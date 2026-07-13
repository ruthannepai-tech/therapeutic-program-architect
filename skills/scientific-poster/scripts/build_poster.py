#!/usr/bin/env python3
"""Build a single-slide, large-format scientific poster (.pptx) from a JSON spec.

Usage:
    python build_poster.py spec.json [--out poster.pptx]

The spec is a JSON object. See SKILL.md for the full schema and the palette
catalog. This script is deliberately self-contained: it also runs a legibility
check and writes `<out>_legibility.txt` and `<out>_speaker_notes.md` next to the
poster so the presenter has everything for both in-person and Zoom delivery.

Requires: python-pptx, pillow, qrcode.
"""
import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Default typeface for all poster text. Arial is near-universal on presentation
# machines and projects cleanly at reading distance.
FONT_NAME = "Arial"

# ---------------------------------------------------------------- palettes ----
# Each palette is chosen for high text/background contrast at reading distance.
PALETTES = {
    "clinical_blue":  dict(primary="#1F3A5F", accent="#2E86C1", panel_bg="#FFFFFF",
                           page_bg="#EAF0F6", text="#1A1A1A", panel_title_text="#FFFFFF"),
    "forest":         dict(primary="#1D4032", accent="#2E8B57", panel_bg="#FFFFFF",
                           page_bg="#EAF3EE", text="#152318", panel_title_text="#FFFFFF"),
    "maroon":         dict(primary="#5A1A2B", accent="#A23A52", panel_bg="#FFFFFF",
                           page_bg="#F4EAED", text="#211015", panel_title_text="#FFFFFF"),
    "slate":          dict(primary="#2C3E50", accent="#D35400", panel_bg="#FFFFFF",
                           page_bg="#ECEFF1", text="#1A1A1A", panel_title_text="#FFFFFF"),
    "teal":           dict(primary="#0B4F6C", accent="#0E8CA8", panel_bg="#FFFFFF",
                           page_bg="#E6F2F5", text="#0A1E26", panel_title_text="#FFFFFF"),
    "plum":           dict(primary="#3B2A50", accent="#8E44AD", panel_bg="#FFFFFF",
                           page_bg="#EFEAF4", text="#1E1428", panel_title_text="#FFFFFF"),
    "charcoal_gold":  dict(primary="#22252A", accent="#977914", panel_bg="#FFFFFF",
                           page_bg="#EDEEF0", text="#1A1A1A", panel_title_text="#FFFFFF"),
    "dark":           dict(primary="#0E1116", accent="#4C9BE0", panel_bg="#1B222C",
                           page_bg="#0E1116", text="#EAEEF3", panel_title_text="#FFFFFF"),
    # sophisticated two-accent palettes tuned for the magazine template
    "editorial_teal": dict(primary="#122A4A", accent="#0E7C7B", accent2="#C25E1A",
                           panel_bg="#FFFFFF", page_bg="#EEF2F5", text="#16202B",
                           panel_title_text="#FFFFFF", subtitle_text="#CFE3E3"),
    "midnight_coral": dict(primary="#1A1F3A", accent="#5A5DD8", accent2="#D64545",
                           panel_bg="#FFFFFF", page_bg="#EDEEF4", text="#171A24",
                           panel_title_text="#FFFFFF", subtitle_text="#C9CBEC"),
    "sage_gold":      dict(primary="#22322A", accent="#3F7D62", accent2="#9A7A12",
                           panel_bg="#FFFFFF", page_bg="#EEF2EE", text="#15201A",
                           panel_title_text="#FFFFFF", subtitle_text="#CFE0D5"),
}


# ------------------------------------------------------------- small helpers ---
def _hex(c):
    return RGBColor.from_string(c.lstrip("#").upper())


def _rel_lum(hexstr):
    """WCAG relative luminance for a hex colour."""
    h = hexstr.lstrip("#")
    r, g, b = (int(h[i:i + 2], 16) / 255 for i in (0, 2, 4))
    def _lin(c):
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)


def _contrast(c1, c2):
    l1, l2 = _rel_lum(c1), _rel_lum(c2)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


def _fit(path, box_w_in, box_h_in):
    """Largest (w, h) in inches that fits `path` inside the box, aspect preserved."""
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = box_w_in / box_h_in
    if ar > box_ar:
        return box_w_in, box_w_in / ar
    return box_h_in * ar, box_h_in


def _panel(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _hex(fill)
    if line:
        sh.line.color.rgb = _hex(line); sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _textbox(slide, l, t, w, h, paragraphs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of paragraphs; each paragraph = list of (text, pt, hexcolor, bold)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    first = True
    for para in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(6)
        for run in para:
            txt, size, color, bold = run[0], run[1], run[2], run[3]
            italic = run[4] if len(run) > 4 else False
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.name = FONT_NAME
            r.font.color.rgb = _hex(color)
    return tb


# ------------------------------------------------------------- legibility ------
def legibility_report(spec, pal, warnings):
    """Return a human-readable legibility/accessibility report string."""
    lines = ["LEGIBILITY & ACCESSIBILITY CHECK", "=" * 34, ""]
    body_pt = spec.get("body_pt", 32)
    head_pt = spec.get("head_pt", 44)
    title_pt = spec.get("title_pt", 80)
    W = spec.get("width_in", 48)

    lines.append("Font sizes (large-format reading-distance floors):")
    def _flag(label, pt, floor):
        ok = "OK " if pt >= floor else ">> "
        return f"  {ok}{label}: {pt}pt (floor {floor}pt)"
    lines.append(_flag("Title", title_pt, 60))
    lines.append(_flag("Section headings", head_pt, 36))
    lines.append(_flag("Body text", body_pt, 24))
    lines.append("")

    # contrast checks
    lines.append("Contrast ratios (WCAG AA large text needs >= 3.0; body >= 4.5):")
    if spec.get("template") == "magazine":
        # magazine headings are white text on filled accent/accent2 strips
        wt = pal.get("panel_title_text", "#FFFFFF")
        pairs = [
            ("Body text on panel",      pal["text"], pal["panel_bg"], 4.5),
            ("Heading strip (accent)",  wt, pal["accent"], 3.0),
            ("Heading strip (accent2)", wt, pal.get("accent2", pal["accent"]), 3.0),
            ("Title on bar",            wt, pal["primary"], 4.5),
        ]
    else:
        pairs = [
            ("Body text on panel", pal["text"], pal["panel_bg"], 4.5),
            ("Heading on panel",   pal["accent"], pal["panel_bg"], 3.0),
            ("Title on bar",       pal["panel_title_text"], pal["primary"], 4.5),
        ]
    for label, fg, bg, floor in pairs:
        cr = _contrast(fg, bg)
        ok = "OK " if cr >= floor else ">> "
        lines.append(f"  {ok}{label}: {cr:.1f}:1 (need {floor}:1)")
    lines.append("")

    if warnings:
        lines.append("Layout warnings:")
        for w in warnings:
            lines.append(f"  >> {w}")
    else:
        lines.append("Layout warnings: none")
    lines.append("")
    lines.append(f"Poster size: {W} x {spec.get('height_in', 36)} in "
                 f"({spec.get('columns', 3)} columns)")
    return "\n".join(lines)


# ------------------------------------------------------------- main builder ----
PPTX_MAX_IN = 56.0  # PowerPoint hard cap on slide width/height


def _fit_title_pt(title, box_w_in, want_pt, floor_pt=60, lines=2, char_w=0.52):
    """Shrink the title font so it fits box_w_in across `lines` lines, but never
    below floor_pt (the large-format legibility minimum). Returns a pt size."""
    cap_per_line = max(1, int((box_w_in * lines) / (char_w * want_pt / 72)))
    if len(title) <= cap_per_line:
        return want_pt
    # scale down proportionally to the overflow, clamp at the floor
    scaled = want_pt * cap_per_line / len(title)
    return max(floor_pt, round(scaled))


def _clamp_dims(W, H, warnings):
    """PowerPoint refuses slides > 56in on a side. Scale the pair down, keeping
    aspect ratio, so a spec that asks for e.g. 60in doesn't silently crash."""
    longest = max(W, H)
    if longest > PPTX_MAX_IN:
        scale = PPTX_MAX_IN / longest
        warnings.append(
            f"Requested {W:.1f}x{H:.1f}in exceeds PowerPoint's {PPTX_MAX_IN:.0f}in cap; "
            f"scaled to {W*scale:.1f}x{H*scale:.1f}in (same aspect ratio). "
            f"Print shops can enlarge the PDF; on-screen it is unaffected.")
        return W * scale, H * scale
    return W, H


def _resolve_palette(spec):
    if isinstance(spec.get("palette"), str):
        pal = dict(PALETTES.get(spec["palette"], PALETTES["clinical_blue"]))
    elif isinstance(spec.get("palette"), dict):
        pal = dict(PALETTES["clinical_blue"]); pal.update(spec["palette"])
    else:
        pal = dict(PALETTES["clinical_blue"])
    pal.setdefault("accent2", pal["accent"])
    pal.setdefault("subtitle_text", pal["panel_title_text"])
    return pal


def build_poster(spec, out="poster.pptx"):
    # Two templates: "classic" (default, simple multi-column) and "magazine"
    # (16:9-friendly, splashy — subtitle band, key-stats ribbon, header strips,
    # figure captions, callouts). Pick with spec["template"].
    if spec.get("template") == "magazine":
        return _build_magazine(spec, out)

    W, H = spec.get("width_in", 48), spec.get("height_in", 36)
    pal = _resolve_palette(spec)

    warnings = []
    W, H = _clamp_dims(W, H, warnings)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _hex(pal["page_bg"])

    M = spec.get("margin_in", 1.0)
    tb_h = spec.get("title_h_in", 4.5)
    _panel(slide, 0, 0, W, tb_h, pal["primary"])

    logo_w = 0
    logo = spec.get("logo")
    if logo and os.path.exists(logo):
        lw, lh = _fit(logo, 3.5, tb_h - 1.2)
        slide.shapes.add_picture(logo, Inches(M), Inches((tb_h - lh) / 2),
                                 Inches(lw), Inches(lh))
        logo_w = lw + 0.5

    _title_w = W - 2 * M - logo_w - 0.5
    _title_pt = spec.get("title_pt") or _fit_title_pt(spec["title"], _title_w, 80, lines=2)
    title_runs = [[(spec["title"], _title_pt, pal["panel_title_text"], True)]]
    if spec.get("authors"):
        title_runs.append([(spec["authors"], spec.get("author_pt", 40),
                            pal["panel_title_text"], False)])
    if spec.get("affiliation"):
        title_runs.append([(spec["affiliation"], spec.get("affil_pt", 30),
                            pal["panel_title_text"], False)])
    _textbox(slide, M + logo_w, 0.2, W - 2 * M - logo_w - 0.5, tb_h - 0.4,
             title_runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    ncol = spec.get("columns", 3)
    gutter = spec.get("gutter_in", 0.8)
    foot_h = spec.get("footer_h_in", 2.5)
    body_top = tb_h + 0.6
    body_h = H - body_top - foot_h - 0.4
    col_w = (W - 2 * M - (ncol - 1) * gutter) / ncol

    body_pt = spec.get("body_pt", 32)
    head_pt = spec.get("head_pt", 44)

    cols = [[] for _ in range(ncol)]
    for i, sec in enumerate(spec["sections"]):
        c = sec.get("column", i % ncol)
        cols[min(c, ncol - 1)].append(sec)

    max_img_h = spec.get("max_img_h_in", 6.5)

    def _measure(sec):
        """Estimate a section's content heights: (head_h, text_h, img_h, natural_panel_h)."""
        title = sec.get("title", "")
        body = sec.get("body", "")
        img = sec.get("image")
        head_h = 0.9 if title else 0
        if body:
            cpl = max(10, int(col_w * 1.9 * (32 / body_pt)))
            lines = 0
            for para in body.split("\n"):
                lines += max(1, int(len(para) / cpl) + 1)
            text_h = lines * (body_pt / 72 * 1.25) + 0.4
        else:
            text_h = 0
        img_h = 0
        if img and os.path.exists(img):
            _, ih = _fit(img, col_w - 0.6, max_img_h)
            img_h = ih + 0.3
            if sec.get("caption"):
                cap_pt = spec.get("caption_pt", 18)
                cpl_c = max(10, int(col_w * 1.9 * (32 / cap_pt)))
                clines = max(1, int(len(sec["caption"]) / cpl_c) + 1)
                img_h += clines * (cap_pt / 72 * 1.2) + 0.15
        elif img:
            warnings.append(f"Image not found: {img}")
        natural = head_h + text_h + img_h + 0.5
        return head_h, text_h, img_h, natural

    for ci in range(ncol):
        x = M + ci * (col_w + gutter)
        secs = cols[ci]
        if not secs:
            continue

        measured = [_measure(s) for s in secs]
        naturals = [m[3] for m in measured]
        total_natural = sum(naturals)
        n = len(secs)

        # Vertical justification: distribute the body height across panels + gaps so the
        # column fills the canvas instead of stacking at the top and leaving dead space.
        # Slack (positive) is shared: part grows panels proportionally, part widens gaps.
        min_gap = 0.5
        # The placement loop below draws (n + 1) gaps: a leading gap before the first
        # panel plus one after each panel. Reserve all of them here so the justified
        # column fills exactly body_h and never overshoots the bottom edge.
        slack = body_h - total_natural - (n + 1) * min_gap
        if slack < -0.01:
            # Content genuinely exceeds the column: keep min gaps, let it run and warn.
            warnings.append(
                f"Column {ci + 1} content is taller than the body area; text/figures may be "
                f"crowded. Shorten text, move a section to another column, reduce max_img_h_in, "
                f"or use a taller poster.")
            gap = min_gap
            panel_hs = naturals
        else:
            # Grow panels to absorb up to ~55% of slack (padding inside panels reads as
            # deliberate whitespace); the rest widens the gaps between panels. Cap any single
            # panel's growth so a lone short section doesn't balloon absurdly.
            # Panels hug their content; most slack becomes clean inter-panel spacing.
            # Figure-bearing panels get first claim on growth (their figure fills it);
            # text-only panels grow little so they don't leave a void under the text.
            grow_budget = slack * 0.35
            gap_budget = slack - grow_budget
            gap = min_gap + (gap_budget / (n + 1))  # also pad top/bottom via leading gap
            has_img = [bool(s.get("image") and os.path.exists(s.get("image", ""))) for s in secs]
            weight = [(2.2 if hi else 1.0) * nat for hi, (_, _, _, nat) in zip(has_img, measured)]
            tw = sum(weight) or 1.0
            panel_hs = []
            for hi, (_, _, _, nat), wt in zip(has_img, measured, weight):
                share = (wt / tw) * grow_budget
                cap = nat * (1.4 if hi else 0.35)  # text-only panels grow modestly
                panel_hs.append(nat + min(share, cap))

        y = body_top + (gap if slack >= -0.01 else 0)
        for sec, (head_h, text_h, img_h, _nat), panel_h in zip(secs, measured, panel_hs):
            title = sec.get("title", "")
            body = sec.get("body", "")
            img = sec.get("image")

            _panel(slide, x, y, col_w, panel_h, pal["panel_bg"], line=pal["accent"])
            iy = y + 0.25
            if title:
                _textbox(slide, x + 0.1, iy, col_w - 0.2, head_h,
                         [[(title, head_pt, pal["accent"], True)]],
                         anchor=MSO_ANCHOR.MIDDLE)
                iy += head_h + 0.05
            if body:
                paras = [[(ln, body_pt, pal["text"], False)] for ln in body.split("\n")]
                # keep body directly under the header (no floating gap); any panel
                # slack falls to the bottom, then figures fill it or it becomes even padding
                _textbox(slide, x + 0.1, iy, col_w - 0.2, text_h, paras)
                iy += text_h
            if img and os.path.exists(img):
                iw, ih = _fit(img, col_w - 0.6, max_img_h)
                cap = sec.get("caption", "")
                cap_pt = spec.get("caption_pt", 18)
                cap_h = 0.0
                if cap:
                    cpl_c = max(10, int(col_w * 1.9 * (32 / cap_pt)))
                    clines = max(1, int(len(cap) / cpl_c) + 1)
                    cap_h = clines * (cap_pt / 72 * 1.2) + 0.15
                # centre the figure (+ caption block) in the panel's remaining space
                avail = y + panel_h - iy - 0.2
                off = max(0, (avail - ih - cap_h) / 2)
                slide.shapes.add_picture(img, Inches(x + (col_w - iw) / 2), Inches(iy + off),
                                         Inches(iw), Inches(ih))
                if cap:
                    _textbox(slide, x + 0.2, iy + off + ih + 0.05, col_w - 0.4, cap_h,
                             [[(cap, cap_pt, pal["text"], False, True)]])
            y += panel_h + gap

    # footer + QR
    _panel(slide, 0, H - foot_h, W, foot_h, pal["primary"])
    qr_w = 0
    if spec.get("qr_url"):
        try:
            import qrcode
            qpath = os.path.splitext(out)[0] + "_qr.png"
            qrcode.make(spec["qr_url"]).save(qpath)
            qs = foot_h - 0.6
            slide.shapes.add_picture(qpath, Inches(W - M - qs), Inches(H - foot_h + 0.3),
                                     Inches(qs), Inches(qs))
            qr_w = qs + 0.4
        except Exception as e:
            warnings.append(f"QR code skipped: {e}")
    foot_runs = [[(spec.get("footer_text", ""), spec.get("footer_pt", 26),
                   pal["panel_title_text"], False)]]
    _textbox(slide, M, H - foot_h, W - 2 * M - qr_w, foot_h, foot_runs,
             anchor=MSO_ANCHOR.MIDDLE)

    # speaker notes
    notes_text = spec.get("speaker_notes", "")
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    prs.save(out)

    # sidecar deliverables (report the sizes actually used, incl. auto-fit title)
    eff = dict(spec)
    eff["title_pt"] = _title_pt
    eff["width_in"], eff["height_in"] = W, H
    report = legibility_report(eff, pal, warnings)
    with open(os.path.splitext(out)[0] + "_legibility.txt", "w") as f:
        f.write(report)
    if notes_text:
        with open(os.path.splitext(out)[0] + "_speaker_notes.md", "w") as f:
            f.write("# Speaker notes / talking points\n\n" + notes_text + "\n")

    return out, warnings, report


def _write_sidecars(out, spec, pal, warnings, notes_text):
    report = legibility_report(spec, pal, warnings)
    with open(os.path.splitext(out)[0] + "_legibility.txt", "w") as f:
        f.write(report)
    if notes_text:
        with open(os.path.splitext(out)[0] + "_speaker_notes.md", "w") as f:
            f.write("# Speaker notes / talking points\n\n" + notes_text + "\n")
    return report


def _header_strip(slide, x, y, w, h, fill, title, pt, text_color):
    """A filled rounded header bar with white section-title text."""
    _panel(slide, x, y, w, h, fill)
    _textbox(slide, x + 0.1, y, w - 0.2, h,
             [[(title, pt, text_color, True)]], anchor=MSO_ANCHOR.MIDDLE)


def _build_magazine(spec, out="poster.pptx"):
    """Splashy, content-dense magazine template. 16:9-friendly.

    Extra spec fields beyond the classic schema:
      template: "magazine"
      subtitle: str            -- one-line takeaway, shown in a band under the title
      stats: [{value,label}]   -- key-numbers ribbon under the title band
      sections[].caption: str  -- italic caption printed under a figure
      sections[].callout: str  -- text rendered in an accent2-tinted highlight box
      sections[].accent: "accent"|"accent2"  -- header-strip colour (round-robin if unset)
    """
    W, H = spec.get("width_in", 48), spec.get("height_in", 27)  # 16:9 default
    pal = _resolve_palette(spec)
    warnings = []
    W, H = _clamp_dims(W, H, warnings)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _hex(pal["page_bg"])

    M = spec.get("margin_in", 0.9)
    # ---- title band with a thin accent2 rule at the bottom ----
    tb_h = spec.get("title_h_in", 3.4)
    _panel(slide, 0, 0, W, tb_h, pal["primary"])
    _panel(slide, 0, tb_h - 0.12, W, 0.12, pal["accent2"])  # accent rule

    logo_w = 0
    logo = spec.get("logo")
    if logo and os.path.exists(logo):
        lw, lh = _fit(logo, 3.2, tb_h - 1.0)
        slide.shapes.add_picture(logo, Inches(M), Inches((tb_h - lh) / 2),
                                 Inches(lw), Inches(lh))
        logo_w = lw + 0.4

    title_pt = spec.get("title_pt") or _fit_title_pt(
        spec["title"], W - 2 * M - logo_w, 62, lines=2)
    sub_pt = spec.get("subtitle_pt", 30)
    title_runs = [[(spec["title"], title_pt, pal["panel_title_text"], True)]]
    if spec.get("subtitle"):
        title_runs.append([(spec["subtitle"], sub_pt, pal["subtitle_text"], False)])
    meta = " · ".join(x for x in [spec.get("authors"), spec.get("affiliation")] if x)
    if meta:
        title_runs.append([(meta, spec.get("author_pt", 22), pal["subtitle_text"], False)])
    _textbox(slide, M + logo_w, 0.15, W - 2 * M - logo_w, tb_h - 0.3,
             title_runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    y_cursor = tb_h + 0.35

    # ---- key-stats ribbon ----
    stats = spec.get("stats", [])
    if stats:
        rib_h = spec.get("stats_h_in", 2.2)
        n = len(stats)
        gap = 0.35
        card_w = (W - 2 * M - (n - 1) * gap) / n
        for i, st in enumerate(stats):
            cx = M + i * (card_w + gap)
            _panel(slide, cx, y_cursor, card_w, rib_h, pal["panel_bg"], line=pal["accent"])
            _panel(slide, cx, y_cursor, 0.16, rib_h, pal["accent2"])  # accent tab
            _textbox(slide, cx + 0.25, y_cursor + 0.12, card_w - 0.35, rib_h - 0.2,
                     [[(str(st["value"]), spec.get("stat_value_pt", 40), pal["accent"], True)],
                      [(st["label"], spec.get("stat_label_pt", 18), pal["text"], False)]],
                     anchor=MSO_ANCHOR.MIDDLE)
        y_cursor += rib_h + 0.5

    # ---- body columns ----
    ncol = spec.get("columns", 4)
    gutter = spec.get("gutter_in", 0.6)
    foot_h = spec.get("footer_h_in", 1.9)
    body_top = y_cursor
    body_h = H - body_top - foot_h - 0.35
    col_w = (W - 2 * M - (ncol - 1) * gutter) / ncol

    body_pt = spec.get("body_pt", 24)
    head_pt = spec.get("head_pt", 36)  # meets the large-format heading floor
    head_h = spec.get("head_strip_h_in", 0.95)
    max_img_h = spec.get("max_img_h_in", 5.0)
    cap_pt = spec.get("caption_pt", 16)

    cols = [[] for _ in range(ncol)]
    for i, sec in enumerate(spec["sections"]):
        c = sec.get("column", i % ncol)
        cols[min(c, ncol - 1)].append(sec)

    def _measure(sec):
        body = sec.get("body", "")
        img = sec.get("image")
        cap = sec.get("caption", "")
        callout = sec.get("callout", "")
        h = head_h + 0.1
        if body:
            cpl = max(8, int(col_w * 1.9 * (24 / body_pt)))
            lines = sum(max(1, int(len(p) / cpl) + 1) for p in body.split("\n"))
            h += lines * (body_pt / 72 * 1.22) + 0.3
        if img and os.path.exists(img):
            _, ih = _fit(img, col_w - 0.4, max_img_h)
            h += ih + 0.25
            if cap:
                cpl = max(8, int(col_w * 1.9 * (24 / cap_pt)))
                clines = sum(max(1, int(len(p) / cpl) + 1) for p in cap.split("\n"))
                h += clines * (cap_pt / 72 * 1.2) + 0.15
        elif img:
            warnings.append(f"Image not found: {img}")
        if callout:
            cpl = max(8, int(col_w * 1.9 * (24 / body_pt)))
            clines = sum(max(1, int(len(p) / cpl) + 1) for p in callout.split("\n"))
            h += clines * (body_pt / 72 * 1.22) + 0.5
        return h + 0.35

    ci_accent = 0
    for ci in range(ncol):
        x = M + ci * (col_w + gutter)
        secs = cols[ci]
        if not secs:
            continue
        naturals = [_measure(s) for s in secs]
        total = sum(naturals)
        n = len(secs)
        min_gap = 0.4
        # placement draws a leading gap + one gap after each panel = (n + 1) gaps
        slack = body_h - total - (n + 1) * min_gap
        if slack < -0.01:
            warnings.append(f"Column {ci + 1} content exceeds body height; shorten text, "
                            f"move a section, reduce max_img_h_in, or add a column.")
            gap = min_gap; heights = naturals; lead = 0
        else:
            grow = slack * 0.5
            gap_budget = slack - grow
            gap = min_gap + gap_budget / (n + 1)
            lead = gap
            heights = [nat + min((nat / total) * grow if total else grow / n, nat * 0.8)
                       for nat in naturals]
        y = body_top + lead
        for sec, ph in zip(secs, heights):
            title = sec.get("title", "")
            body = sec.get("body", "")
            img = sec.get("image")
            cap = sec.get("caption", "")
            callout = sec.get("callout", "")
            acc_key = sec.get("accent") or ("accent2" if ci_accent % 2 else "accent")
            ci_accent += 1
            acc = pal.get(acc_key, pal["accent"])

            _panel(slide, x, y, col_w, ph, pal["panel_bg"], line=pal["accent"])
            iy = y + 0.12
            if title:
                _header_strip(slide, x + 0.12, iy, col_w - 0.24, head_h, acc,
                              title, head_pt, pal["panel_title_text"])
                iy += head_h + 0.15
            if body:
                paras = [[(ln, body_pt, pal["text"], False)] for ln in body.split("\n")]
                bh = ph - (iy - y) - 0.2
                _textbox(slide, x + 0.14, iy, col_w - 0.28, bh, paras)
                # advance by measured text height, not full remaining
                cpl = max(8, int(col_w * 1.9 * (24 / body_pt)))
                lines = sum(max(1, int(len(p) / cpl) + 1) for p in body.split("\n"))
                iy += lines * (body_pt / 72 * 1.22) + 0.25
            if img and os.path.exists(img):
                iw, ih = _fit(img, col_w - 0.4, max_img_h)
                slide.shapes.add_picture(img, Inches(x + (col_w - iw) / 2), Inches(iy),
                                         Inches(iw), Inches(ih))
                iy += ih + 0.1
                if cap:
                    _textbox(slide, x + 0.14, iy, col_w - 0.28, 0.9,
                             [[(cap, cap_pt, pal["text"], False)]])
                    iy += 0.5
            if callout:
                # tinted highlight box using a light wash of accent2
                co_h = y + ph - iy - 0.15
                co_h = max(co_h, 0.6)
                cbox = _panel(slide, x + 0.14, iy, col_w - 0.28, co_h, pal["accent2"])
                _textbox(slide, x + 0.24, iy, col_w - 0.48, co_h,
                         [[(callout, body_pt, pal["panel_title_text"], True)]],
                         anchor=MSO_ANCHOR.MIDDLE)
            y += ph + gap

    # ---- footer + QR ----
    _panel(slide, 0, H - foot_h, W, foot_h, pal["primary"])
    _panel(slide, 0, H - foot_h, W, 0.1, pal["accent2"])
    qr_w = 0
    if spec.get("qr_url"):
        try:
            import qrcode
            qpath = os.path.splitext(out)[0] + "_qr.png"
            qrcode.make(spec["qr_url"]).save(qpath)
            qs = foot_h - 0.5
            slide.shapes.add_picture(qpath, Inches(W - M - qs), Inches(H - foot_h + 0.28),
                                     Inches(qs), Inches(qs))
            qr_w = qs + 0.4
        except Exception as e:
            warnings.append(f"QR code skipped: {e}")
    _textbox(slide, M, H - foot_h, W - 2 * M - qr_w, foot_h,
             [[(spec.get("footer_text", ""), spec.get("footer_pt", 22),
                pal["panel_title_text"], False)]], anchor=MSO_ANCHOR.MIDDLE)

    notes_text = spec.get("speaker_notes", "")
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    prs.save(out)
    # report against the sizes actually used by the magazine template (its defaults
    # differ from classic), so the legibility check reflects the rendered poster.
    eff = dict(spec)
    eff.setdefault("title_pt", title_pt)
    eff.setdefault("head_pt", head_pt)
    eff.setdefault("body_pt", body_pt)
    eff["width_in"], eff["height_in"] = W, H
    report = _write_sidecars(out, eff, pal, warnings, notes_text)
    return out, warnings, report


def main():
    ap = argparse.ArgumentParser(description="Build a scientific poster (.pptx) from a JSON spec.")
    ap.add_argument("spec", help="Path to the poster spec JSON file.")
    ap.add_argument("--out", default="poster.pptx", help="Output .pptx path.")
    args = ap.parse_args()

    with open(args.spec) as f:
        spec = json.load(f)
    out, warnings, report = build_poster(spec, args.out)
    print(f"Saved {out}")
    print(report)
    if warnings:
        print("\nWARNINGS:")
        for w in warnings:
            print(" -", w)
    else:
        print("\nNo warnings.")


if __name__ == "__main__":
    main()
